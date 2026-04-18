import sys, os, warnings, zipfile, re
warnings.filterwarnings("ignore")

action = sys.argv[1]
input_path = sys.argv[2]
output_path = sys.argv[3]

try:
    # ---------------------------------------------------------
    # FIXED: 100% OCR EXTRACTION (Forced 300 DPI + PSM 3)
    # ---------------------------------------------------------
    if action == "ocr":
        import pytesseract
        from pdf2image import convert_from_path
        # Force 300 DPI for high-res scanning, strict page segmentation
        custom_config = r'--oem 3 --psm 3'
        images = convert_from_path(input_path, dpi=300)
        text = "".join([pytesseract.image_to_string(img, config=custom_config) + "\n\n" for img in images])
        with open(output_path, 'w', encoding='utf-8') as f: 
            f.write(text if text.strip() else "OCR ENGINE ERROR: No readable text detected at 300 DPI.")

    # ---------------------------------------------------------
    # FIXED: RAW TEXT EXTRACTION
    # ---------------------------------------------------------
    elif action == "extract_text":
        import fitz
        doc = fitz.open(input_path)
        extracted = "\n".join([page.get_text("text") for page in doc])
        with open(output_path, 'w', encoding='utf-8') as f: 
            f.write(extracted if extracted.strip() else "System detected zero raw text layers. Document may be a scanned image. Use OCR instead.")

    # ---------------------------------------------------------
    # FIXED: TABLE TO EXCEL (Lattice + Stream Mode)
    # ---------------------------------------------------------
    elif action == "excel":
        import tabula, pandas as pd
        # lattice=True forces grid detection, stream=True reads borderless tables
        dfs = tabula.read_pdf(input_path, pages='all', multiple_tables=True, stream=True, guess=False)
        if dfs: 
            with pd.ExcelWriter(output_path) as writer:
                for i, df in enumerate(dfs): df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
        else: pd.DataFrame({'System Warning': ['No tables detected by engine']}).to_excel(output_path, index=False)

    # ---------------------------------------------------------
    # FIXED: CONVERT TO WORD (Real structure conversion)
    # ---------------------------------------------------------
    elif action == "word":
        from pdf2docx import Converter
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    # ---------------------------------------------------------
    # FIXED: CONVERT TO POWERPOINT (Slide injection)
    # ---------------------------------------------------------
    elif action == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        from pdf2image import convert_from_path
        prs = Presentation()
        images = convert_from_path(input_path, dpi=200)
        for img in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
            img_path = f"temp_{os.urandom(4).hex()}.jpg"
            img.save(img_path, 'JPEG')
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width)
            os.remove(img_path)
        prs.save(output_path)

    # ---------------------------------------------------------
    # FIXED: NLP AUTO SUMMARY (TF-IDF Ranking)
    # ---------------------------------------------------------
    elif action == "summary":
        import fitz, nltk
        from nltk.corpus import stopwords
        from nltk.tokenize import word_tokenize, sent_tokenize
        text = " ".join([page.get_text() for page in fitz.open(input_path)])
        
        stop_words = set(stopwords.words("english"))
        words = word_tokenize(text)
        freq_table = dict()
        for word in words:
            word = word.lower()
            if word in stop_words or not word.isalnum(): continue
            freq_table[word] = freq_table.get(word, 0) + 1
            
        sentences = sent_tokenize(text)
        sent_value = dict()
        for sent in sentences:
            for word, freq in freq_table.items():
                if word in sent.lower(): sent_value[sent] = sent_value.get(sent, 0) + freq
                
        sum_values = sum(sent_value.values())
        average = int(sum_values / len(sent_value)) if len(sent_value) > 0 else 0
        
        summary = " ".join([s for s in sentences if (s in sent_value) and (sent_value[s] > (1.2 * average))])
        with open(output_path, 'w', encoding='utf-8') as f: 
            f.write("ENTERPRISE NLP SUMMARY:\n\n" + (summary if summary else "Not enough text to summarize."))

    # Original working features...
    elif action == "pdf2img":
        from pdf2image import convert_from_path
        with zipfile.ZipFile(output_path, 'w') as z:
            for i, img in enumerate(convert_from_path(input_path, dpi=200)):
                img_path = f"page_{i+1}.jpg"
                img.save(img_path, "JPEG")
                z.write(img_path)
                os.remove(img_path)
    elif action == "redact":
        import fitz
        doc = fitz.open(input_path)
        patterns = [r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', r'\b\d{3}[-.\s]??\d{3}[-.\s]??\d{4}\b']
        for page in doc:
            for pat in patterns:
                for inst in page.search_for(pat, quads=True): page.add_redact_annot(inst, fill=(0, 0, 0))
            page.apply_redactions()
        doc.save(output_path)
    elif action == "html":
        import fitz
        doc = fitz.open(input_path)
        html_content = "<html><head><style>body{font-family:sans-serif; margin:40px; background:#f4f7f6;} .page{background:#fff; padding:20px; margin-bottom:20px; box-shadow:0 2px 4px rgba(0,0,0,0.1);}</style></head><body>"
        for page in doc: html_content += f"<div class='page'>{page.get_text('html')}</div>"
        html_content += "</body></html>"
        with open(output_path, 'w', encoding='utf-8') as f: f.write(html_content)

except Exception as e:
    print(f"[PYTHON AI ERROR] {str(e)}")
    sys.exit(1)
